#!/usr/bin/env python
# -*- coding: utf-8 -*-
import os, sys, io
from pptx import Presentation
from pptx.util import Emu

src = r"e:\CA001\Infomat\docs\外部参考\数据标准管理培训.pptx"
out = r"e:\CA001\Infomat\_tmp\data_std_pptx_raw.txt"

prs = Presentation(src)

buf = io.StringIO()

def w(s=""):
    buf.write(s)
    buf.write("\n")

w(f"=== PPTX: {os.path.basename(src)} ===")
w(f"Slide width  : {prs.slide_width} EMU  ({Emu(prs.slide_width).inches:.2f} in)")
w(f"Slide height : {prs.slide_height} EMU  ({Emu(prs.slide_height).inches:.2f} in)")
w(f"Total slides : {len(prs.slides)}")
w(f"Slide masters: {len(prs.slide_masters)}")
w(f"Slide layouts: {len(prs.slide_layouts)}")
w()

for idx, slide in enumerate(prs.slides, 1):
    w(f"{'='*78}")
    w(f"Slide {idx}  |  layout: {slide.slide_layout.name}")
    w(f"{'='*78}")
    for shp_i, shape in enumerate(slide.shapes, 1):
        kind = shape.shape_type
        name = shape.name
        is_ph = shape.is_placeholder
        ph_type = None
        ph_idx = None
        if is_ph:
            try:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
            except Exception:
                pass
        text_runs = []
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                runs = []
                for run in para.runs:
                    runs.append(run.text)
                joined = "".join(runs)
                if joined.strip() or any(r.strip() for r in runs):
                    text_runs.append({
                        "level": para.level,
                        "alignment": str(para.alignment) if para.alignment is not None else None,
                        "text": joined,
                        "runs": runs,
                    })
        table_data = None
        if shape.has_table:
            tbl = shape.table
            rows = []
            for r in tbl.rows:
                row_cells = []
                for c in r.cells:
                    row_cells.append(c.text)
                rows.append(row_cells)
            table_data = rows
        pic_info = None
        if shape.shape_type == 13:
            try:
                pic_info = {
                    "filename": shape.image.filename if hasattr(shape, "image") else None,
                    "content_type": shape.image.content_type if hasattr(shape, "image") else None,
                    "ext": shape.image.ext if hasattr(shape, "image") else None,
                }
            except Exception as e:
                pic_info = f"err: {e}"
        chart_info = None
        if shape.has_chart:
            try:
                ch = shape.chart
                chart_info = {
                    "chart_type": str(ch.chart_type),
                    "has_title": bool(ch.has_title),
                    "title": ch.chart_title.text_frame.text if ch.has_title else None,
                    "series_count": len(ch.series),
                }
            except Exception as e:
                chart_info = f"err: {e}"
        group_info = None
        if shape.shape_type == 6:
            group_info = f"shapes_in_group={len(shape.shapes)}"

        geom = {
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
        }
        w(f"")
        w(f"  [Shape {shp_i}] name={name!r} type={kind} placeholder={is_ph}"
          + (f" ph_type={ph_type} ph_idx={ph_idx}" if is_ph else ""))
        w(f"    geom(EMU): {geom}")
        if text_runs:
            for tr in text_runs:
                lvl = "  " * tr["level"]
                w(f"    TEXT[L{tr['level']}]: {lvl}{tr['text']!r}")
        if table_data:
            w(f"    TABLE rows={len(table_data)} cols={len(table_data[0]) if table_data else 0}")
            for ri, row in enumerate(table_data, 1):
                for ci, cell in enumerate(row, 1):
                    preview = cell.replace("\n", " | ")
                    if len(preview) > 200:
                        preview = preview[:200] + "…"
                    w(f"      R{ri}C{ci}: {preview!r}")
        if pic_info:
            w(f"    PICTURE: {pic_info}")
        if chart_info:
            w(f"    CHART: {chart_info}")
        if group_info:
            w(f"    GROUP: {group_info}")
        if not (text_runs or table_data or pic_info or chart_info or group_info):
            try:
                alt = shape.element.attrib
                w(f"    (no text/table/picture) attrs={ {k:v for k,v in alt.items() if 'alt' in k.lower()} }")
            except Exception:
                pass

w("=== DONE ===")

with open(out, "w", encoding="utf-8") as f:
    f.write(buf.getvalue())
print(f"Written {len(buf.getvalue())} chars to {out}")
